#!/usr/bin/env python3
"""
File Converter Module - Converts various file formats to text
"""
import os
import pandas as pd
import PyPDF2
from docx import Document
import chardet
import tempfile
import hashlib
from datetime import datetime
import re
import json
import xml.etree.ElementTree as ET
import zipfile
import email
from email import policy
from bs4 import BeautifulSoup
import markdown
import yaml

class FileConverter:
    def __init__(self):
        self.max_file_size = 50 * 1024 * 1024  # 50MB
        self.allowed_extensions = {
            '.txt', '.csv', '.xlsx', '.xls', '.pdf', '.docx', '.doc',
            '.json', '.xml', '.html', '.htm', '.md', '.yaml', '.yml',
            '.log', '.conf', '.ini', '.cfg', '.properties',
            '.py', '.js', '.java', '.cpp', '.c', '.h', '.hpp',
            '.sh', '.bat', '.ps1', '.sql', '.css', '.scss',
            '.eml', '.msg', '.rtf', '.odt', '.pptx', '.ppt'
        }
    
    def is_allowed_file(self, filename):
        """Check if file extension is allowed"""
        ext = os.path.splitext(filename)[1].lower()
        return ext in self.allowed_extensions
    
    def sanitize_filename(self, filename):
        """Sanitize filename for security"""
        # Remove path components
        filename = os.path.basename(filename)
        # Remove special characters
        filename = re.sub(r'[^\w\s.-]', '', filename)
        # Add timestamp hash
        timestamp = datetime.now().strftime('%Y%m%d%H%M%S')
        hash_suffix = hashlib.md5(f"{filename}{timestamp}".encode()).hexdigest()[:8]
        name, ext = os.path.splitext(filename)
        return f"{name}_{hash_suffix}{ext}"
    
    def detect_encoding(self, file_path):
        """Detect file encoding"""
        with open(file_path, 'rb') as f:
            raw_data = f.read(10000)  # Read first 10KB
            result = chardet.detect(raw_data)
            return result['encoding'] or 'utf-8'
    
    def convert_to_text(self, file_path, filename):
        """Convert file to text based on extension"""
        ext = os.path.splitext(filename)[1].lower()
        
        try:
            if ext in ['.txt', '.log', '.conf', '.ini', '.cfg', '.properties']:
                return self.convert_text_file(file_path)
            elif ext in ['.csv']:
                return self.convert_csv(file_path)
            elif ext in ['.xlsx', '.xls']:
                return self.convert_excel(file_path)
            elif ext == '.pdf':
                return self.convert_pdf(file_path)
            elif ext in ['.docx', '.doc']:
                return self.convert_word(file_path)
            elif ext == '.json':
                return self.convert_json(file_path)
            elif ext in ['.xml', '.html', '.htm']:
                return self.convert_markup(file_path, ext)
            elif ext in ['.md']:
                return self.convert_markdown(file_path)
            elif ext in ['.yaml', '.yml']:
                return self.convert_yaml(file_path)
            elif ext in ['.py', '.js', '.java', '.cpp', '.c', '.h', '.hpp', '.sh', '.bat', '.ps1', '.sql', '.css', '.scss']:
                return self.convert_code_file(file_path, ext)
            elif ext in ['.eml', '.msg']:
                return self.convert_email(file_path)
            elif ext == '.rtf':
                return self.convert_rtf(file_path)
            elif ext == '.odt':
                return self.convert_odt(file_path)
            elif ext in ['.pptx', '.ppt']:
                return self.convert_powerpoint(file_path)
            else:
                # Try to read as text with encoding detection
                return self.convert_text_file(file_path)
        except Exception as e:
            return f"ファイルの変換中にエラーが発生しました: {str(e)}"
    
    def convert_text_file(self, file_path):
        """Convert text file with encoding detection"""
        encoding = self.detect_encoding(file_path)
        try:
            with open(file_path, 'r', encoding=encoding) as f:
                return f.read()
        except UnicodeDecodeError:
            # Fallback to binary read and decode with errors='ignore'
            with open(file_path, 'rb') as f:
                content = f.read()
                return content.decode('utf-8', errors='ignore')
    
    def convert_csv(self, file_path):
        """Convert CSV to text"""
        encoding = self.detect_encoding(file_path)
        try:
            df = pd.read_csv(file_path, encoding=encoding)
            return df.to_string(index=False)
        except:
            # Try with different encodings
            for enc in ['utf-8', 'shift-jis', 'cp932', 'euc-jp', 'iso-2022-jp']:
                try:
                    df = pd.read_csv(file_path, encoding=enc)
                    return df.to_string(index=False)
                except:
                    continue
            return "CSVファイルの読み込みに失敗しました"
    
    def convert_excel(self, file_path):
        """Convert Excel to text"""
        try:
            # Read all sheets
            xl_file = pd.ExcelFile(file_path)
            text_output = []
            
            for sheet_name in xl_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                text_output.append(f"\n=== シート: {sheet_name} ===\n")
                text_output.append(df.to_string(index=False))
            
            return '\n'.join(text_output)
        except Exception as e:
            return f"Excelファイルの読み込みエラー: {str(e)}"
    
    def convert_pdf(self, file_path):
        """Convert PDF to text"""
        try:
            text_output = []
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    text_output.append(f"\n--- ページ {page_num} ---\n")
                    text_output.append(page.extract_text())
            return '\n'.join(text_output)
        except Exception as e:
            return f"PDFファイルの読み込みエラー: {str(e)}"
    
    def convert_word(self, file_path):
        """Convert Word document to text"""
        try:
            doc = Document(file_path)
            text_output = []
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_output.append(para.text)
            
            # Extract tables
            for table in doc.tables:
                text_output.append("\n[テーブル]")
                for row in table.rows:
                    row_text = '\t'.join(cell.text for cell in row.cells)
                    text_output.append(row_text)
            
            return '\n'.join(text_output)
        except Exception as e:
            return f"Wordファイルの読み込みエラー: {str(e)}"
    
    def convert_json(self, file_path):
        """Convert JSON to formatted text"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            return json.dumps(data, ensure_ascii=False, indent=2)
        except Exception as e:
            return f"JSONファイルの読み込みエラー: {str(e)}"
    
    def convert_markup(self, file_path, ext):
        """Convert HTML/XML to text"""
        try:
            encoding = self.detect_encoding(file_path)
            with open(file_path, 'r', encoding=encoding) as f:
                content = f.read()
            
            if ext in ['.html', '.htm']:
                soup = BeautifulSoup(content, 'html.parser')
                # Remove script and style elements
                for script in soup(["script", "style"]):
                    script.decompose()
                return soup.get_text(separator='\n', strip=True)
            else:  # XML
                root = ET.fromstring(content)
                return self._xml_to_text(root)
        except Exception as e:
            return f"マークアップファイルの読み込みエラー: {str(e)}"
    
    def _xml_to_text(self, element, level=0):
        """Recursively convert XML element to text"""
        indent = "  " * level
        text = f"{indent}<{element.tag}>"
        
        if element.text and element.text.strip():
            text += f" {element.text.strip()}"
        
        text += "\n"
        
        for child in element:
            text += self._xml_to_text(child, level + 1)
        
        if element.tail and element.tail.strip():
            text += f"{indent}{element.tail.strip()}\n"
        
        return text
    
    def convert_markdown(self, file_path):
        """Convert Markdown to plain text"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                md_content = f.read()
            # Convert to HTML then extract text
            html = markdown.markdown(md_content)
            soup = BeautifulSoup(html, 'html.parser')
            return soup.get_text(separator='\n', strip=True)
        except Exception as e:
            return self.convert_text_file(file_path)
    
    def convert_yaml(self, file_path):
        """Convert YAML to formatted text"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = yaml.safe_load(f)
            return yaml.dump(data, default_flow_style=False, allow_unicode=True)
        except Exception as e:
            return self.convert_text_file(file_path)
    
    def convert_code_file(self, file_path, ext):
        """Convert code files with syntax information"""
        content = self.convert_text_file(file_path)
        language_map = {
            '.py': 'Python',
            '.js': 'JavaScript',
            '.java': 'Java',
            '.cpp': 'C++',
            '.c': 'C',
            '.h': 'C Header',
            '.hpp': 'C++ Header',
            '.sh': 'Shell Script',
            '.bat': 'Batch Script',
            '.ps1': 'PowerShell',
            '.sql': 'SQL',
            '.css': 'CSS',
            '.scss': 'SCSS'
        }
        language = language_map.get(ext, 'Code')
        return f"=== {language}ファイル ===\n{content}"
    
    def convert_email(self, file_path):
        """Convert email files to text"""
        try:
            with open(file_path, 'rb') as f:
                msg = email.message_from_binary_file(f, policy=policy.default)
            
            text_output = []
            text_output.append(f"差出人: {msg.get('From', 'N/A')}")
            text_output.append(f"宛先: {msg.get('To', 'N/A')}")
            text_output.append(f"件名: {msg.get('Subject', 'N/A')}")
            text_output.append(f"日付: {msg.get('Date', 'N/A')}")
            text_output.append("\n--- 本文 ---\n")
            
            # Get email body
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        text_output.append(part.get_payload(decode=True).decode('utf-8', errors='ignore'))
            else:
                text_output.append(msg.get_payload(decode=True).decode('utf-8', errors='ignore'))
            
            return '\n'.join(text_output)
        except Exception as e:
            return f"メールファイルの読み込みエラー: {str(e)}"
    
    def convert_rtf(self, file_path):
        """Convert RTF to text (basic implementation)"""
        try:
            with open(file_path, 'rb') as f:
                content = f.read()
            # Basic RTF to text conversion (removes RTF control words)
            text = content.decode('utf-8', errors='ignore')
            # Remove RTF control words
            text = re.sub(r'\\[a-z]+\d*\s?', '', text)
            text = re.sub(r'[{}]', '', text)
            return text.strip()
        except Exception as e:
            return f"RTFファイルの読み込みエラー: {str(e)}"
    
    def convert_odt(self, file_path):
        """Convert ODT (OpenDocument Text) to text"""
        try:
            text_output = []
            with zipfile.ZipFile(file_path, 'r') as z:
                # Read content.xml
                with z.open('content.xml') as f:
                    content = f.read()
                    root = ET.fromstring(content)
                    # Extract text from all text elements
                    for elem in root.iter():
                        if elem.text:
                            text_output.append(elem.text)
            return ' '.join(text_output)
        except Exception as e:
            return f"ODTファイルの読み込みエラー: {str(e)}"
    
    def convert_powerpoint(self, file_path):
        """Convert PowerPoint to text"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_output = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text_output.append(f"\n=== スライド {slide_num} ===\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_output.append(shape.text)
            
            return '\n'.join(text_output)
        except ImportError:
            return "PowerPointファイルの変換にはpython-pptxが必要です"
        except Exception as e:
            return f"PowerPointファイルの読み込みエラー: {str(e)}"